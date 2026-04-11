import PyPDF2
import docx
from pptx import Presentation
from typing import Optional

def extract_text_from_file(file_path: str) -> Optional[str]:
    """
    Extract text content from a PDF, DOCX, or PPTX file.
    
    Args:
        file_path: Path to the file
        
    Returns:
        Extracted text as string, or None if extraction fails
    """
    try:
        lower_path = file_path.lower()
        if lower_path.endswith('.pdf'):
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()
                
        elif lower_path.endswith('.docx'):
            doc = docx.Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
            
        elif lower_path.endswith('.pptx'):
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
            
        else:
            print(f"Unsupported file format: {file_path}")
            return None
            
    except Exception as e:
        print(f"Error extracting file text: {e}")
        return None
